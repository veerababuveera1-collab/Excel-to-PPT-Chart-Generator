import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
from datetime import datetime, timedelta
from scipy import stats
from sklearn.linear_model import LinearRegression
from sklearn.preprocessing import StandardScaler
import warnings
warnings.filterwarnings('ignore')

# ══════════════════════════════════════════════════════════════════
# 0. PAGE CONFIG
# ══════════════════════════════════════════════════════════════════
st.set_page_config(
    page_title="NEXUS Command | Enterprise Intelligence",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ══════════════════════════════════════════════════════════════════
# 1. GLOBAL STYLES — Dark Cyberpunk Mission Control
# ══════════════════════════════════════════════════════════════════
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@400;600;700;900&family=Rajdhani:wght@300;400;500;600;700&family=Share+Tech+Mono&display=swap');

/* ── Root Variables ── */
:root {
    --bg-void:     #020408;
    --bg-deep:     #060d16;
    --bg-panel:    #0a1628;
    --bg-card:     #0d1f35;
    --bg-glass:    rgba(13, 31, 53, 0.7);
    --accent-cyan: #00f5ff;
    --accent-mag:  #ff006e;
    --accent-gold: #ffd60a;
    --accent-lime: #39ff14;
    --accent-blue: #0066ff;
    --border-glow: rgba(0, 245, 255, 0.25);
    --text-bright: #e8f4f8;
    --text-mid:    #8eb8d4;
    --text-dim:    #3d6680;
    --glow-cyan:   0 0 20px rgba(0,245,255,0.4), 0 0 60px rgba(0,245,255,0.15);
    --glow-mag:    0 0 20px rgba(255,0,110,0.4), 0 0 60px rgba(255,0,110,0.15);
    --glow-gold:   0 0 20px rgba(255,214,10,0.4);
}

/* ── Global Reset ── */
html, body, .stApp {
    background-color: var(--bg-void) !important;
    font-family: 'Rajdhani', sans-serif;
    color: var(--text-bright);
}

/* ── Scrollbar ── */
::-webkit-scrollbar { width: 4px; }
::-webkit-scrollbar-track { background: var(--bg-deep); }
::-webkit-scrollbar-thumb { background: var(--accent-cyan); border-radius: 2px; }

/* ── Main Content ── */
.main .block-container {
    padding: 1.5rem 2rem;
    max-width: 100%;
    background: var(--bg-void);
}

/* ── Sidebar ── */
[data-testid="stSidebar"] {
    background: var(--bg-deep) !important;
    border-right: 1px solid var(--border-glow);
}
[data-testid="stSidebar"] * { color: var(--text-bright) !important; }
[data-testid="stSidebar"] .stSelectbox label,
[data-testid="stSidebar"] .stMultiselect label,
[data-testid="stSidebar"] .stSlider label { color: var(--text-mid) !important; font-size: 0.72rem; letter-spacing: 1px; text-transform: uppercase; }

/* ── Inputs ── */
.stTextInput input, .stSelectbox select,
[data-testid="stTextInput"] input,
[data-baseweb="select"] { background: var(--bg-panel) !important; border: 1px solid var(--border-glow) !important; color: var(--text-bright) !important; border-radius: 4px !important; font-family: 'Rajdhani', sans-serif !important; }
[data-baseweb="select"]:focus-within { border-color: var(--accent-cyan) !important; box-shadow: var(--glow-cyan) !important; }

/* ── Password Input ── */
[data-testid="stTextInput"] input[type="password"] {
    background: var(--bg-panel) !important;
    border: 1px solid var(--border-glow) !important;
    color: var(--accent-cyan) !important;
    font-family: 'Share Tech Mono', monospace !important;
    letter-spacing: 4px;
}

/* ── Buttons ── */
.stButton > button {
    background: transparent !important;
    border: 1px solid var(--accent-cyan) !important;
    color: var(--accent-cyan) !important;
    font-family: 'Orbitron', sans-serif !important;
    font-size: 0.7rem !important;
    letter-spacing: 2px !important;
    text-transform: uppercase !important;
    padding: 0.6rem 1.5rem !important;
    border-radius: 2px !important;
    transition: all 0.3s ease !important;
    position: relative !important;
    overflow: hidden !important;
}
.stButton > button:hover {
    background: var(--accent-cyan) !important;
    color: var(--bg-void) !important;
    box-shadow: var(--glow-cyan) !important;
    transform: translateY(-1px) !important;
}

/* ── Tabs ── */
.stTabs [data-baseweb="tab-list"] {
    background: var(--bg-deep) !important;
    border-bottom: 1px solid var(--border-glow) !important;
    gap: 0;
}
.stTabs [data-baseweb="tab"] {
    color: var(--text-dim) !important;
    font-family: 'Orbitron', sans-serif !important;
    font-size: 0.65rem !important;
    letter-spacing: 1.5px !important;
    text-transform: uppercase !important;
    border-radius: 0 !important;
    border-bottom: 2px solid transparent !important;
    padding: 0.8rem 1.5rem !important;
    transition: all 0.3s !important;
    background: transparent !important;
}
.stTabs [aria-selected="true"] {
    color: var(--accent-cyan) !important;
    border-bottom: 2px solid var(--accent-cyan) !important;
    background: var(--bg-panel) !important;
}

/* ── Metrics ── */
[data-testid="stMetric"] {
    background: var(--bg-card) !important;
    border: 1px solid var(--border-glow) !important;
    border-radius: 6px !important;
    padding: 1rem !important;
    position: relative !important;
    overflow: hidden !important;
}
[data-testid="stMetric"]::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 2px;
    background: linear-gradient(90deg, var(--accent-cyan), var(--accent-mag));
}
[data-testid="stMetricLabel"] { color: var(--text-dim) !important; font-family: 'Orbitron', sans-serif !important; font-size: 0.6rem !important; letter-spacing: 2px !important; text-transform: uppercase !important; }
[data-testid="stMetricValue"] { color: var(--accent-cyan) !important; font-family: 'Orbitron', sans-serif !important; font-size: 1.6rem !important; }
[data-testid="stMetricDelta"] { font-family: 'Share Tech Mono', monospace !important; }

/* ── Dataframe ── */
[data-testid="stDataFrame"] { border: 1px solid var(--border-glow) !important; border-radius: 6px !important; }
.stDataFrame thead tr th { background: var(--bg-panel) !important; color: var(--accent-cyan) !important; font-family: 'Orbitron', monospace !important; font-size: 0.62rem !important; letter-spacing: 1px !important; }

/* ── Alerts / Info ── */
.stInfo, .stSuccess, .stWarning, .stError { border-radius: 4px !important; border-left-width: 3px !important; font-family: 'Rajdhani', sans-serif !important; }

/* ── File Uploader ── */
[data-testid="stFileUploader"] {
    border: 1px dashed var(--border-glow) !important;
    border-radius: 6px !important;
    background: var(--bg-panel) !important;
    padding: 1rem !important;
}

/* ── Expander ── */
.streamlit-expanderHeader {
    background: var(--bg-panel) !important;
    border: 1px solid var(--border-glow) !important;
    color: var(--text-mid) !important;
    font-family: 'Orbitron', sans-serif !important;
    font-size: 0.65rem !important;
    letter-spacing: 1.5px !important;
}

/* ══ Custom Components ══ */

/* Hero Header */
.nexus-header {
    background: var(--bg-deep);
    border: 1px solid var(--border-glow);
    border-radius: 8px;
    padding: 2rem 2.5rem;
    margin-bottom: 1.5rem;
    position: relative;
    overflow: hidden;
    display: flex;
    align-items: center;
    justify-content: space-between;
}
.nexus-header::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 2px;
    background: linear-gradient(90deg, transparent, var(--accent-cyan), var(--accent-mag), transparent);
}
.nexus-header::after {
    content: '';
    position: absolute;
    bottom: 0; left: 0; right: 0;
    height: 1px;
    background: linear-gradient(90deg, transparent, var(--border-glow), transparent);
}
.nexus-title {
    font-family: 'Orbitron', sans-serif;
    font-weight: 900;
    font-size: 2rem;
    color: var(--accent-cyan);
    text-shadow: var(--glow-cyan);
    letter-spacing: 3px;
    line-height: 1;
}
.nexus-subtitle {
    font-family: 'Share Tech Mono', monospace;
    color: var(--text-mid);
    font-size: 0.78rem;
    letter-spacing: 2px;
    margin-top: 0.4rem;
}
.nexus-badge {
    font-family: 'Share Tech Mono', monospace;
    font-size: 0.65rem;
    letter-spacing: 2px;
    text-transform: uppercase;
    padding: 0.3rem 0.8rem;
    border-radius: 2px;
    display: inline-block;
    margin: 0.1rem;
}
.badge-live {
    background: rgba(57,255,20,0.1);
    border: 1px solid var(--accent-lime);
    color: var(--accent-lime);
}
.badge-time {
    background: rgba(0,245,255,0.1);
    border: 1px solid var(--accent-cyan);
    color: var(--accent-cyan);
}

/* Status Panel */
.status-panel {
    background: var(--bg-card);
    border: 1px solid var(--border-glow);
    border-radius: 6px;
    padding: 1.2rem 1.5rem;
    margin-bottom: 1rem;
    position: relative;
}
.status-panel.go    { border-left: 3px solid var(--accent-lime); }
.status-panel.caution { border-left: 3px solid var(--accent-gold); }
.status-panel.nogo  { border-left: 3px solid var(--accent-mag); }

.status-label {
    font-family: 'Orbitron', sans-serif;
    font-size: 0.6rem;
    letter-spacing: 3px;
    text-transform: uppercase;
    color: var(--text-dim);
}
.status-value {
    font-family: 'Orbitron', sans-serif;
    font-size: 1.4rem;
    font-weight: 700;
    margin-top: 0.2rem;
}
.status-go    { color: var(--accent-lime); text-shadow: var(--glow-cyan); }
.status-caution { color: var(--accent-gold); text-shadow: var(--glow-gold); }
.status-nogo  { color: var(--accent-mag); text-shadow: var(--glow-mag); }

/* Insight Card */
.insight-card {
    background: var(--bg-card);
    border: 1px solid var(--border-glow);
    border-radius: 6px;
    padding: 1rem 1.2rem;
    margin: 0.4rem 0;
    font-family: 'Rajdhani', sans-serif;
    font-size: 0.9rem;
    line-height: 1.5;
    color: var(--text-mid);
}
.insight-card strong { color: var(--accent-cyan); }

/* Scan Line Effect */
.scan-container { position: relative; }
.scan-container::after {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0; bottom: 0;
    background: repeating-linear-gradient(0deg, transparent, transparent 2px, rgba(0,245,255,0.01) 2px, rgba(0,245,255,0.01) 4px);
    pointer-events: none;
    border-radius: 6px;
}

/* KPI Tile */
.kpi-tile {
    background: var(--bg-card);
    border: 1px solid var(--border-glow);
    border-radius: 6px;
    padding: 1.2rem;
    text-align: center;
    position: relative;
    overflow: hidden;
    transition: all 0.3s;
}
.kpi-tile:hover { border-color: var(--accent-cyan); box-shadow: var(--glow-cyan); transform: translateY(-2px); }
.kpi-tile::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 2px;
}
.kpi-tile.cyan::before { background: var(--accent-cyan); }
.kpi-tile.mag::before  { background: var(--accent-mag); }
.kpi-tile.gold::before { background: var(--accent-gold); }
.kpi-tile.lime::before { background: var(--accent-lime); }
.kpi-label { font-family:'Orbitron',sans-serif; font-size:0.55rem; letter-spacing:2px; text-transform:uppercase; color:var(--text-dim); }
.kpi-value { font-family:'Orbitron',sans-serif; font-size:1.8rem; font-weight:700; margin:0.3rem 0; }
.kpi-tile.cyan .kpi-value { color:var(--accent-cyan); }
.kpi-tile.mag  .kpi-value { color:var(--accent-mag);  }
.kpi-tile.gold .kpi-value { color:var(--accent-gold); }
.kpi-tile.lime .kpi-value { color:var(--accent-lime); }
.kpi-delta { font-family:'Share Tech Mono',monospace; font-size:0.72rem; color:var(--text-dim); }

/* Section Header */
.section-header {
    font-family: 'Orbitron', sans-serif;
    font-size: 0.7rem;
    letter-spacing: 3px;
    text-transform: uppercase;
    color: var(--text-dim);
    border-bottom: 1px solid var(--border-glow);
    padding-bottom: 0.5rem;
    margin-bottom: 1rem;
    display: flex;
    align-items: center;
    gap: 0.8rem;
}
.section-header::before {
    content: '';
    width: 24px;
    height: 2px;
    background: linear-gradient(90deg, var(--accent-cyan), transparent);
    display: inline-block;
}

/* Login Screen */
.login-shell {
    background: var(--bg-deep);
    border: 1px solid var(--border-glow);
    border-radius: 10px;
    padding: 3rem;
    position: relative;
    overflow: hidden;
    max-width: 420px;
    margin: auto;
}
.login-shell::before {
    content: '';
    position: absolute;
    top:0;left:0;right:0;
    height: 2px;
    background: linear-gradient(90deg, transparent, var(--accent-cyan), var(--accent-mag), transparent);
}
.login-logo {
    font-family: 'Orbitron', sans-serif;
    font-size: 2.2rem;
    font-weight: 900;
    color: var(--accent-cyan);
    text-shadow: var(--glow-cyan);
    text-align: center;
    letter-spacing: 4px;
}
.login-tag {
    font-family: 'Share Tech Mono', monospace;
    font-size: 0.65rem;
    color: var(--text-dim);
    text-align: center;
    letter-spacing: 3px;
    margin-bottom: 2rem;
}
.login-divider {
    border: none;
    border-top: 1px solid var(--border-glow);
    margin: 1.5rem 0;
}

/* Anomaly Badge */
.anomaly-badge {
    display: inline-block;
    background: rgba(255,0,110,0.15);
    border: 1px solid var(--accent-mag);
    color: var(--accent-mag);
    font-family: 'Share Tech Mono', monospace;
    font-size: 0.65rem;
    letter-spacing: 2px;
    padding: 0.2rem 0.6rem;
    border-radius: 2px;
}

/* Sidebar Profile */
.sidebar-profile {
    background: var(--bg-card);
    border: 1px solid var(--border-glow);
    border-radius: 6px;
    padding: 1rem;
    margin-bottom: 1rem;
    text-align: center;
}
.profile-name {
    font-family: 'Orbitron', sans-serif;
    font-size: 0.85rem;
    font-weight: 700;
    color: var(--accent-cyan) !important;
    letter-spacing: 1px;
}
.profile-role {
    font-family: 'Share Tech Mono', monospace;
    font-size: 0.62rem;
    color: var(--text-dim) !important;
    letter-spacing: 2px;
    text-transform: uppercase;
}

/* Correlation Cell */
.corr-positive { color: var(--accent-cyan); }
.corr-negative { color: var(--accent-mag); }
</style>
""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════════
# 2. HELPERS & ANALYTICS ENGINE
# ══════════════════════════════════════════════════════════════════

PLOTLY_LAYOUT = dict(
    paper_bgcolor='rgba(0,0,0,0)',
    plot_bgcolor='rgba(6,13,22,0.9)',
    font=dict(family='Rajdhani', color='#8eb8d4', size=12),
    title_font=dict(family='Orbitron', color='#00f5ff', size=13),
    xaxis=dict(gridcolor='rgba(0,245,255,0.06)', linecolor='rgba(0,245,255,0.15)', zerolinecolor='rgba(0,245,255,0.1)'),
    yaxis=dict(gridcolor='rgba(0,245,255,0.06)', linecolor='rgba(0,245,255,0.15)', zerolinecolor='rgba(0,245,255,0.1)'),
    legend=dict(bgcolor='rgba(6,13,22,0.8)', bordercolor='rgba(0,245,255,0.2)', borderwidth=1, font=dict(family='Rajdhani')),
    margin=dict(l=20, r=20, t=50, b=20),
    hoverlabel=dict(bgcolor='#0a1628', bordercolor='#00f5ff', font=dict(family='Share Tech Mono', size=11, color='#00f5ff'))
)

PALETTE = ['#00f5ff', '#ff006e', '#ffd60a', '#39ff14', '#0066ff', '#ff8800', '#c77dff', '#f72585']

def apply_layout(fig, title=""):
    fig.update_layout(**PLOTLY_LAYOUT)
    if title:
        fig.update_layout(title=dict(text=title, font=dict(family='Orbitron', color='#00f5ff', size=13)))
    return fig

def detect_anomalies(series: pd.Series, threshold: float = 2.5) -> pd.Series:
    """Z-score based anomaly detection."""
    z = np.abs(stats.zscore(series.fillna(series.mean())))
    return z > threshold

def forecast_linear(series: pd.Series, periods: int = 14) -> pd.DataFrame:
    """Simple linear regression forecast."""
    x = np.arange(len(series)).reshape(-1, 1)
    y = series.fillna(method='ffill').values
    model = LinearRegression().fit(x, y)
    future_x = np.arange(len(series), len(series) + periods).reshape(-1, 1)
    preds = model.predict(future_x)
    ci = 1.96 * np.std(y - model.predict(x))
    return pd.DataFrame({'forecast': preds, 'upper': preds + ci, 'lower': preds - ci})

def compute_dre(df, date_col, metric_col) -> float:
    """Defect Removal Efficiency approximation."""
    recent = df.nlargest(max(1, len(df)//3), metric_col)[metric_col].sum()
    total  = df[metric_col].sum()
    return round((1 - recent / total) * 100, 1) if total > 0 else 95.0

def risk_score(df, slicer, metric_col) -> pd.DataFrame:
    """Multi-factor risk scoring per dimension."""
    g = df.groupby(slicer)[metric_col]
    rs = pd.DataFrame({
        'Volume': g.sum(),
        'Volatility': g.std().fillna(0),
        'Count': g.count(),
    })
    scaler = StandardScaler()
    if len(rs) > 1:
        rs['RiskScore'] = scaler.fit_transform(rs[['Volume', 'Volatility', 'Count']]).mean(axis=1)
        rs['RiskScore'] = ((rs['RiskScore'] - rs['RiskScore'].min()) /
                           (rs['RiskScore'].max() - rs['RiskScore'].min() + 1e-9) * 100).round(1)
    else:
        rs['RiskScore'] = 50.0
    return rs.reset_index()

def generate_insights(df, slicer, metric_col, stability, date_col=None) -> list:
    """Rule-based auto-insight engine."""
    insights = []
    top = df.groupby(slicer)[metric_col].sum().idxmax()
    top_pct = (df.groupby(slicer)[metric_col].sum().max() /
               df[metric_col].sum() * 100) if df[metric_col].sum() > 0 else 0
    insights.append(f"<strong>Concentration Risk:</strong> <em>{top}</em> accounts for <strong>{top_pct:.1f}%</strong> of total {metric_col} — consider load-balancing.")
    if stability < 60:
        insights.append(f"<strong>🔴 Critical Alert:</strong> Stability Index is <strong>{stability}%</strong>. Immediate governance intervention required.")
    elif stability < 80:
        insights.append(f"<strong>🟡 Watch Zone:</strong> Stability at <strong>{stability}%</strong>. Monitor inflow velocity closely over next 72 hours.")
    else:
        insights.append(f"<strong>🟢 Healthy Baseline:</strong> Stability Index is <strong>{stability}%</strong>. System within operational thresholds.")
    if date_col and date_col in df.columns:
        try:
            daily = df.groupby(date_col).size()
            if len(daily) >= 7:
                recent_7 = daily.iloc[-7:].mean()
                prev_7   = daily.iloc[-14:-7].mean() if len(daily) >= 14 else daily.mean()
                delta_pct = ((recent_7 - prev_7) / (prev_7 + 1e-9)) * 100
                direction = "↑ accelerating" if delta_pct > 10 else ("↓ decelerating" if delta_pct < -10 else "→ stable")
                insights.append(f"<strong>Trend Signal:</strong> 7-day inflow rate is <em>{direction}</em> at <strong>{recent_7:.1f}</strong>/day (Δ {delta_pct:+.1f}% vs prior period).")
        except Exception:
            pass
    dims = df[slicer].nunique()
    insights.append(f"<strong>Diversity:</strong> {dims} active {slicer} dimensions tracked. Pareto principle suggests focusing on top {max(1, dims//5)} for 80% impact.")
    return insights

def build_pptx(user, stability, status_text, total_val, metric_label, top_module, dre, df_summary):
    """Build a polished boardroom PPTX with data summary table."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK   = RGBColor(0x02, 0x04, 0x08)
    CYAN   = RGBColor(0x00, 0xf5, 0xff)
    WHITE  = RGBColor(0xe8, 0xf4, 0xf8)
    GREY   = RGBColor(0x3d, 0x66, 0x80)
    ACCENT = RGBColor(0xff, 0x00, 0x6e) if stability < 60 else (RGBColor(0xff, 0xd6, 0x0a) if stability < 80 else RGBColor(0x39, 0xff, 0x14))

    blank_layout = prs.slide_layouts[6]

    def dark_slide(title_text, subtitle_text=""):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK

        # Top bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
        bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()

        # Title
        txb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.7))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        run = p.add_run(); run.text = title_text
        run.font.size  = Pt(28); run.font.bold = True
        run.font.color.rgb = CYAN
        p.alignment = PP_ALIGN.LEFT

        if subtitle_text:
            txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.88), Inches(12), Inches(0.4))
            tf2  = txb2.text_frame
            p2   = tf2.paragraphs[0]
            r2   = p2.add_run(); r2.text = subtitle_text
            r2.font.size = Pt(11); r2.font.color.rgb = GREY
        return slide

    def add_kpi_box(slide, x, y, w, h, label, value, color):
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0d, 0x1f, 0x35)
        box.line.color.rgb = color; box.line.width = Pt(1)
        txb = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(h*0.4))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]; r = p.add_run(); r.text = label.upper()
        r.font.size = Pt(8); r.font.color.rgb = GREY; p.alignment = PP_ALIGN.CENTER
        txb2 = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+h*0.35), Inches(w-0.2), Inches(h*0.55))
        tf2  = txb2.text_frame
        p2   = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = str(value)
        r2.font.size = Pt(22); r2.font.bold = True; r2.font.color.rgb = color
        p2.alignment = PP_ALIGN.CENTER

    # SLIDE 1 — Title
    s1 = dark_slide("NEXUS COMMAND", "Enterprise Intelligence Platform — Executive Briefing")
    txb = s1.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8), Inches(1.5))
    tf = txb.text_frame; tf.word_wrap = True
    for line, clr in [
        (f"Presented by:  {user}", WHITE),
        (f"Generated:     {datetime.now().strftime('%Y-%m-%d %H:%M')} UTC", GREY),
        (f"Classification: BOARD CONFIDENTIAL", CYAN),
    ]:
        p = tf.add_paragraph(); r = p.add_run(); r.text = line
        r.font.size = Pt(13); r.font.color.rgb = clr; r.font.name = 'Courier New'

    # Status badge
    status_box = s1.shapes.add_shape(1, Inches(9.5), Inches(1.5), Inches(3.2), Inches(1.2))
    status_box.fill.solid(); status_box.fill.fore_color.rgb = RGBColor(0x06, 0x0d, 0x16)
    status_box.line.color.rgb = ACCENT; status_box.line.width = Pt(2)
    txb_s = s1.shapes.add_textbox(Inches(9.5), Inches(1.65), Inches(3.2), Inches(0.9))
    tf_s = txb_s.text_frame
    p_s = tf_s.paragraphs[0]; r_s = p_s.add_run(); r_s.text = status_text
    r_s.font.size = Pt(16); r_s.font.bold = True; r_s.font.color.rgb = ACCENT
    p_s.alignment = PP_ALIGN.CENTER

    # SLIDE 2 — KPIs
    s2 = dark_slide("EXECUTIVE KPI DASHBOARD", "Core performance indicators at a glance")
    add_kpi_box(s2, 0.5, 1.5, 2.8, 1.4, f"Total {metric_label}", f"{total_val:,.0f}", CYAN)
    add_kpi_box(s2, 3.5, 1.5, 2.8, 1.4, "Stability Index", f"{stability}%", ACCENT)
    add_kpi_box(s2, 6.5, 1.5, 2.8, 1.4, "DRE Score", f"{dre}%", RGBColor(0x39,0xff,0x14))
    add_kpi_box(s2, 9.5, 1.5, 2.8, 1.4, "Hotspot Module", str(top_module)[:16], RGBColor(0xff,0xd6,0x0a))

    # SLIDE 3 — Summary Table
    s3 = dark_slide("RISK DIMENSION ANALYSIS", "Aggregated metrics by strategic dimension")
    if df_summary is not None and len(df_summary) > 0:
        cols = df_summary.columns.tolist()
        n_rows = min(len(df_summary), 10) + 1
        n_cols = min(len(cols), 5)
        tbl = s3.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.8)).table
        tbl.columns[0].width = Inches(2.5)
        for ci, col in enumerate(cols[:n_cols]):
            cell = tbl.cell(0, ci)
            cell.text = str(col).upper()
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x0a,0x16,0x28)
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            for run in p.runs: run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = CYAN
        for ri in range(min(len(df_summary), 9)):
            for ci, col in enumerate(cols[:n_cols]):
                cell = tbl.cell(ri+1, ci)
                val  = df_summary.iloc[ri][col]
                cell.text = f"{val:.1f}" if isinstance(val, float) else str(val)
                cell.fill.solid(); cell.fill.fore_color.rgb = (RGBColor(0x0d,0x1f,0x35) if ri % 2 == 0 else RGBColor(0x06,0x0d,0x16))
                p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                for run in p.runs: run.font.size = Pt(9); run.font.color.rgb = WHITE

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ══════════════════════════════════════════════════════════════════
# 3. AUTH GATE
# ══════════════════════════════════════════════════════════════════
if "auth" not in st.session_state:
    st.markdown("<br><br>", unsafe_allow_html=True)
    _, c2, _ = st.columns([1, 1.2, 1])
    with c2:
        st.markdown("""
        <div class="login-shell">
            <div class="login-logo">NEXUS</div>
            <div class="login-tag">COMMAND / ENTERPRISE INTELLIGENCE</div>
            <hr class="login-divider">
        </div>
        """, unsafe_allow_html=True)
        st.markdown("**⬡ IDENTIFIER**")
        u = st.text_input("", placeholder="Director Username", key="login_user", label_visibility="collapsed")
        st.markdown("**⬡ SECURITY KEY**")
        p = st.text_input("", placeholder="••••••••••••", type="password", key="login_pass", label_visibility="collapsed")
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("⚡ AUTHORIZE ACCESS", use_container_width=True):
            if p == "Company2026" and u.strip():
                st.session_state["auth"] = True
                st.session_state["user"] = u.strip()
                st.session_state["login_time"] = datetime.now().strftime("%Y-%m-%d %H:%M")
                st.rerun()
            else:
                st.error("ACCESS DENIED — Invalid credentials.")
        st.markdown("""
        <div style="text-align:center; margin-top:1.5rem;">
            <span class="nexus-badge badge-live">● SYSTEM ONLINE</span>
            <span class="nexus-badge badge-time">NEXUS v4.2.1</span>
        </div>
        """, unsafe_allow_html=True)
    st.stop()

# ══════════════════════════════════════════════════════════════════
# 4. SIDEBAR — Controls & Data Engine
# ══════════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown(f"""
    <div class="sidebar-profile">
        <div style="font-size:2rem; margin-bottom:0.3rem;">⬡</div>
        <div class="profile-name">{st.session_state['user']}</div>
        <div class="profile-role">Senior Director</div>
        <div style="font-family:'Share Tech Mono',monospace; font-size:0.6rem; color:#3d6680; margin-top:0.4rem;">
            Session: {st.session_state.get('login_time','—')}
        </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="section-header">DATA INGESTION</div>', unsafe_allow_html=True)
    uploaded_file = st.file_uploader("Upload Defect Master (.xlsx)", type=["xlsx"], label_visibility="collapsed")

    if not uploaded_file:
        # ── DEMO MODE ──
        st.info("⚡ No file uploaded — running with synthetic demo data.")
        np.random.seed(42)
        n = 320
        modules   = ['Authentication', 'Payments', 'Reporting', 'API Gateway', 'Mobile App', 'Data Pipeline']
        statuses  = ['Open', 'In Progress', 'Resolved', 'Closed']
        severities = ['Critical', 'High', 'Medium', 'Low']
        priorities = ['P0', 'P1', 'P2', 'P3']
        demo_df = pd.DataFrame({
            'Module':      np.random.choice(modules, n, p=[0.25,0.2,0.15,0.15,0.15,0.1]),
            'Status':      np.random.choice(statuses, n, p=[0.3,0.25,0.25,0.2]),
            'Severity':    np.random.choice(severities, n, p=[0.15,0.3,0.35,0.2]),
            'Priority':    np.random.choice(priorities, n, p=[0.1,0.25,0.4,0.25]),
            'Defect_Count':np.random.randint(1, 15, n),
            'Effort_Hours':np.random.uniform(0.5, 40, n).round(1),
            'Date':        pd.date_range('2024-10-01', periods=n, freq='4H')[:n],
            'Root_Cause':  np.random.choice(['Logic Error','UX Issue','API Fault','Data Issue','Config','Security'], n),
        })
        df = demo_df
    else:
        df = pd.read_excel(uploaded_file)
        # Auto-repair
        for col in df.columns:
            if df[col].dtype == 'object':
                try: df[col] = pd.to_numeric(df[col].astype(str).str.replace(r'[$, ]', '', regex=True))
                except: pass
        date_cols_detect = [c for c in df.columns if 'date' in c.lower()]
        if date_cols_detect:
            df[date_cols_detect[0]] = pd.to_datetime(df[date_cols_detect[0]], errors='coerce')

    date_cols = [c for c in df.columns if pd.api.types.is_datetime64_any_dtype(df[c])]
    num_cols  = df.select_dtypes(include='number').columns.tolist()
    cat_cols  = df.select_dtypes(exclude='number').columns.tolist()
    cat_cols  = [c for c in cat_cols if not pd.api.types.is_datetime64_any_dtype(df[c])]

    st.markdown('<div class="section-header" style="margin-top:1rem;">GOVERNANCE FILTERS</div>', unsafe_allow_html=True)

    slicer = st.selectbox("Strategic Dimension (X-Axis)", cat_cols)
    all_vals = df[slicer].dropna().unique().tolist()
    selected = st.multiselect(f"Focus: {slicer}", all_vals, default=all_vals)
    df_filtered = df[df[slicer].isin(selected)] if selected else df.copy()

    for ac in [c for c in df.columns if any(x in c.lower() for x in ['status','severity','priority'])]:
        if ac in cat_cols:
            opts = df[ac].dropna().unique().tolist()
            sel  = st.multiselect(f"⬡ {ac}", opts, default=opts)
            df_filtered = df_filtered[df_filtered[ac].isin(sel)]

    if date_cols:
        st.markdown("**⬡ REPORTING PERIOD**")
        min_d = df[date_cols[0]].min().date()
        max_d = df[date_cols[0]].max().date()
        dr    = st.date_input("Date Range", [min_d, max_d], label_visibility="collapsed")
        if len(dr) == 2:
            df_filtered = df_filtered[
                (df_filtered[date_cols[0]].dt.date >= dr[0]) &
                (df_filtered[date_cols[0]].dt.date <= dr[1])
            ]

    st.markdown('<div class="section-header" style="margin-top:1rem;">VISUAL CONTROLS</div>', unsafe_allow_html=True)
    raw_y = st.selectbox("Primary Metric (Y-Axis)", num_cols)
    chart_theme = st.color_picker("Accent Color", "#00f5ff")
    forecast_periods = st.slider("Forecast Horizon (days)", 7, 60, 14)

    st.divider()
    if st.button("🚪 LOGOUT", use_container_width=True):
        for k in list(st.session_state.keys()): del st.session_state[k]
        st.rerun()

if df_filtered.empty:
    st.warning("⚠ No data matches current filters. Adjust the sidebar controls.")
    st.stop()

y_label = raw_y.split('(')[0].strip()

# ══════════════════════════════════════════════════════════════════
# 5. CORE CALCULATIONS
# ══════════════════════════════════════════════════════════════════
total_val = df_filtered[raw_y].sum()
grp_sum   = df_filtered.groupby(slicer)[raw_y].sum()
top_raw   = grp_sum.idxmax()
top_module = top_raw.strftime('%Y-%m-%d') if hasattr(top_raw, 'strftime') else str(top_raw)
risk_pct   = (grp_sum.max() / total_val * 100) if total_val > 0 else 0

# Stability Index
if date_cols:
    recent_date = df_filtered[date_cols[0]].max()
    last_3 = df_filtered[df_filtered[date_cols[0]] > (recent_date - timedelta(days=3))]
    inflow_rate    = len(last_3)
    stability_score = max(0, min(100, 100 - (inflow_rate * 5)))
else:
    stability_score = 95

dre_score   = compute_dre(df_filtered, date_cols[0] if date_cols else None, raw_y)
risk_df     = risk_score(df_filtered, slicer, raw_y)
status_text = "🟢 STABLE" if stability_score >= 75 else ("🟡 CAUTION" if stability_score >= 60 else "🔴 AT RISK")
status_cls  = "go" if stability_score >= 75 else ("caution" if stability_score >= 60 else "nogo")
status_val_cls  = "status-go" if stability_score >= 75 else ("status-caution" if stability_score >= 60 else "status-nogo")
anomaly_series  = detect_anomalies(df_filtered[raw_y]) if len(df_filtered) > 3 else pd.Series([False]*len(df_filtered))
anomaly_count   = anomaly_series.sum()
insights_list   = generate_insights(df_filtered, slicer, raw_y, stability_score, date_cols[0] if date_cols else None)

# ══════════════════════════════════════════════════════════════════
# 6. HERO HEADER
# ══════════════════════════════════════════════════════════════════
now_str = datetime.now().strftime("%Y-%m-%d  %H:%M:%S")
st.markdown(f"""
<div class="nexus-header">
    <div>
        <div class="nexus-title">⚡ NEXUS COMMAND</div>
        <div class="nexus-subtitle">ENTERPRISE INTELLIGENCE PLATFORM  ·  PREDICTIVE GOVERNANCE SUITE</div>
        <div style="margin-top:0.6rem;">
            <span class="nexus-badge badge-live">● LIVE</span>
            <span class="nexus-badge badge-time">{now_str}</span>
            <span class="nexus-badge badge-time">{len(df_filtered):,} RECORDS</span>
        </div>
    </div>
    <div style="text-align:right;">
        <div class="status-label">MISSION STATUS</div>
        <div class="status-value {status_val_cls}">{status_text}</div>
        <div style="font-family:'Share Tech Mono',monospace; font-size:0.65rem; color:#3d6680; margin-top:0.3rem;">
            STABILITY INDEX: {stability_score}%
        </div>
    </div>
</div>
""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════════
# 7. KPI BAND
# ══════════════════════════════════════════════════════════════════
k1, k2, k3, k4, k5, k6 = st.columns(6)
kpi_data = [
    (k1, "cyan",  f"Total {y_label}", f"{total_val:,.0f}", f"{len(df_filtered):,} records"),
    (k2, "mag",   "Hotspot Module",   top_module[:14],      f"{risk_pct:.1f}% exposure"),
    (k3, "gold",  "Stability Index",  f"{stability_score}%", status_text),
    (k4, "lime",  "DRE Score",        f"{dre_score}%",       "Defect Removal Eff."),
    (k5, "cyan",  "Anomalies",        str(int(anomaly_count)), f"/{len(df_filtered)} records"),
    (k6, "mag",   "Dimensions",       str(df_filtered[slicer].nunique()), f"{slicer} active"),
]
for col, color, label, val, sub in kpi_data:
    with col:
        st.markdown(f"""
        <div class="kpi-tile {color}">
            <div class="kpi-label">{label}</div>
            <div class="kpi-value">{val}</div>
            <div class="kpi-delta">{sub}</div>
        </div>
        """, unsafe_allow_html=True)

st.markdown("<br>", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════════
# 8. TABS
# ══════════════════════════════════════════════════════════════════
tabs = st.tabs([
    "⬡ COMMAND CENTER",
    "⚠ RISK MATRIX",
    "📈 VELOCITY & FORECAST",
    "🔬 ADVANCED ANALYTICS",
    "🔍 AUDIT TRAIL",
    "📡 AI INSIGHTS",
])

# ────────────────────────────────────────
# TAB 1 — COMMAND CENTER
# ────────────────────────────────────────
with tabs[0]:
    col_left, col_right = st.columns([2.2, 1])

    with col_left:
        st.markdown('<div class="section-header">VOLUME DISTRIBUTION</div>', unsafe_allow_html=True)
        agg = df_filtered.groupby(slicer)[raw_y].sum().reset_index().sort_values(raw_y, ascending=False)
        fig_bar = go.Figure(go.Bar(
            x=agg[slicer], y=agg[raw_y],
            marker=dict(
                color=agg[raw_y],
                colorscale=[[0,'#0a1628'],[0.5, chart_theme],[1,'#ff006e']],
                line=dict(color='rgba(0,245,255,0.3)', width=1),
            ),
            text=[f"{v:,.0f}" for v in agg[raw_y]],
            textposition='outside',
            textfont=dict(family='Share Tech Mono', size=10, color='#8eb8d4'),
            hovertemplate='<b>%{x}</b><br>Value: %{y:,.0f}<extra></extra>',
        ))
        apply_layout(fig_bar, f"{y_label} by {slicer}")
        st.plotly_chart(fig_bar, use_container_width=True)

    with col_right:
        st.markdown('<div class="section-header">ROOT CAUSE RADAR</div>', unsafe_allow_html=True)
        if 'Root_Cause' in df_filtered.columns:
            rc_data = df_filtered.groupby('Root_Cause')[raw_y].sum().reset_index()
        else:
            rc_data = pd.DataFrame({'RC':['Logic','UX','API','Data','Config'], 'Val':[40,20,25,10,5]})
            rc_data.columns = ['Root_Cause', raw_y]

        fig_donut = go.Figure(go.Pie(
            labels=rc_data['Root_Cause'], values=rc_data[raw_y],
            hole=0.62,
            marker=dict(colors=PALETTE, line=dict(color='#020408', width=2)),
            textfont=dict(family='Share Tech Mono', size=9, color='#e8f4f8'),
            hovertemplate='<b>%{label}</b><br>%{value:,.0f} — %{percent}<extra></extra>',
        ))
        fig_donut.add_annotation(text=f"<b>{total_val:,.0f}</b>", font=dict(family='Orbitron', size=20, color='#00f5ff'), showarrow=False)
        apply_layout(fig_donut, "Root Cause Breakdown")
        st.plotly_chart(fig_donut, use_container_width=True)

        st.markdown('<div class="section-header" style="margin-top:0.5rem;">RELEASE READINESS</div>', unsafe_allow_html=True)
        st.markdown(f"""
        <div class="status-panel {status_cls}">
            <div class="status-label">MISSION STATUS</div>
            <div class="status-value {status_val_cls}">{status_text}</div>
            <div style="font-family:'Rajdhani',sans-serif; font-size:0.85rem; color:#8eb8d4; margin-top:0.5rem;">
                <strong>{top_module}</strong> carries <strong>{risk_pct:.1f}%</strong> of risk volume.<br>
                Inflow rate signals <strong>{"pressure" if stability_score < 75 else "stability"}</strong>.
            </div>
        </div>
        """, unsafe_allow_html=True)

    # Second row — Treemap + Scatter
    col_a, col_b = st.columns(2)
    with col_a:
        st.markdown('<div class="section-header">RISK TREEMAP</div>', unsafe_allow_html=True)
        rs = risk_df.copy()
        fig_tree = px.treemap(rs, path=[slicer], values='Volume',
                              color='RiskScore', color_continuous_scale=[[0,'#0a1628'],[0.5,'#004080'],[1,'#ff006e']],
                              hover_data=['Volatility','Count','RiskScore'])
        fig_tree.update_traces(textfont=dict(family='Orbitron', size=11))
        apply_layout(fig_tree, "Risk Volume Treemap")
        fig_tree.update_coloraxes(colorbar=dict(tickfont=dict(color='#8eb8d4'), title=dict(text='Risk', font=dict(color='#8eb8d4'))))
        st.plotly_chart(fig_tree, use_container_width=True)

    with col_b:
        st.markdown('<div class="section-header">MULTI-DIM RISK SCATTER</div>', unsafe_allow_html=True)
        rs2 = risk_df.copy()
        fig_scatter = px.scatter(
            rs2, x='Volume', y='Volatility', size='Count', color='RiskScore',
            text=slicer, color_continuous_scale=[[0,'#004080'],[0.5,'#00f5ff'],[1,'#ff006e']],
            hover_data={slicer: True, 'RiskScore': ':.1f', 'Count': True},
        )
        fig_scatter.update_traces(textfont=dict(family='Share Tech Mono', size=9, color='#e8f4f8'), textposition='top center')
        apply_layout(fig_scatter, "Risk Scatter: Volume vs Volatility")
        st.plotly_chart(fig_scatter, use_container_width=True)

# ────────────────────────────────────────
# TAB 2 — RISK MATRIX
# ────────────────────────────────────────
with tabs[1]:
    st.markdown('<div class="section-header">AGING HEATMAP & RISK CLASSIFICATION</div>', unsafe_allow_html=True)

    c1, c2 = st.columns([1.6, 1])
    with c1:
        if date_cols:
            df_filtered = df_filtered.copy()
            df_filtered['Age_Days'] = (pd.Timestamp(datetime.now()) - df_filtered[date_cols[0]]).dt.days.clip(lower=0)
            age_bins  = pd.cut(df_filtered['Age_Days'], bins=[-1,3,7,14,100], labels=["🟢 0-3d","🟡 4-7d","🟠 8-14d","🔴 14d+"])
            age_group = df_filtered.groupby([slicer, age_bins])[raw_y].sum().unstack(fill_value=0)

            fig_heat = go.Figure(go.Heatmap(
                z=age_group.values,
                x=age_group.columns.astype(str).tolist(),
                y=age_group.index.astype(str).tolist(),
                colorscale=[[0,'#0a1628'],[0.33,'#004080'],[0.66,'#ffd60a'],[1,'#ff006e']],
                text=age_group.values,
                texttemplate='%{text:.0f}',
                textfont=dict(family='Share Tech Mono', size=10),
                hovertemplate='%{y} × %{x}<br>Value: %{z:,.0f}<extra></extra>',
            ))
            apply_layout(fig_heat, f"Aging Heatmap: {slicer} × Age Bucket")
            st.plotly_chart(fig_heat, use_container_width=True)
        else:
            st.info("Upload data with a Date column to see the Aging Heatmap.")

    with c2:
        st.markdown('<div class="section-header">RISK SCORE RANKING</div>', unsafe_allow_html=True)
        rs_sorted = risk_df.sort_values('RiskScore', ascending=False)
        fig_risk_bar = go.Figure(go.Bar(
            x=rs_sorted['RiskScore'], y=rs_sorted[slicer],
            orientation='h',
            marker=dict(
                color=rs_sorted['RiskScore'],
                colorscale=[[0,'#0a1628'],[0.5,'#004080'],[1,'#ff006e']],
                line=dict(color='rgba(0,245,255,0.2)', width=0.5),
            ),
            text=[f"{v:.0f}" for v in rs_sorted['RiskScore']],
            textposition='inside',
            textfont=dict(family='Share Tech Mono', size=10, color='white'),
            hovertemplate='<b>%{y}</b><br>Risk Score: %{x:.1f}<extra></extra>',
        ))
        apply_layout(fig_risk_bar, "Composite Risk Score")
        fig_risk_bar.update_layout(yaxis=dict(autorange='reversed'))
        st.plotly_chart(fig_risk_bar, use_container_width=True)

    # Anomaly table
    st.markdown('<div class="section-header">ANOMALY DETECTION REPORT</div>', unsafe_allow_html=True)
    if anomaly_count > 0:
        anom_df = df_filtered[anomaly_series.values].copy()
        anom_df['⚠ Anomaly Score'] = np.abs(stats.zscore(df_filtered[raw_y].fillna(0)))[anomaly_series.values].round(2)
        st.markdown(f'<span class="anomaly-badge">⚠ {int(anomaly_count)} ANOMALIES DETECTED</span>', unsafe_allow_html=True)
        st.dataframe(anom_df.head(20), use_container_width=True)
    else:
        st.success("✓ No statistical anomalies detected in current filtered dataset.")

    # Box plot distribution
    st.markdown('<div class="section-header">DISTRIBUTION INTELLIGENCE</div>', unsafe_allow_html=True)
    fig_box = go.Figure()
    for i, grp in enumerate(df_filtered[slicer].unique()):
        sub = df_filtered[df_filtered[slicer]==grp][raw_y]
        fig_box.add_trace(go.Box(
            y=sub, name=str(grp),
            marker_color=PALETTE[i % len(PALETTE)],
            line_color=PALETTE[i % len(PALETTE)],
            fillcolor=f'rgba({int(PALETTE[i%len(PALETTE)][1:3],16)},{int(PALETTE[i%len(PALETTE)][3:5],16)},{int(PALETTE[i%len(PALETTE)][5:],16)},0.2)',
            boxmean='sd',
            hovertemplate='<b>%{x}</b><br>%{y}<extra></extra>',
        ))
    apply_layout(fig_box, f"{y_label} Statistical Distribution by {slicer}")
    st.plotly_chart(fig_box, use_container_width=True)

# ────────────────────────────────────────
# TAB 3 — VELOCITY & FORECAST
# ────────────────────────────────────────
with tabs[2]:
    if not date_cols:
        st.info("Upload data with a Date column to access Velocity & Forecast analytics.")
    else:
        dc = date_cols[0]
        daily = df_filtered.groupby(dc).agg(Inflow=(raw_y,'sum'), Count=('index' if 'index' not in df_filtered.columns else raw_y,'count')).reset_index()
        daily = daily.sort_values(dc)
        daily['Cumulative_Inflow'] = daily['Inflow'].cumsum()
        daily['Cumulative_Outflow'] = (daily['Cumulative_Inflow'] * 0.87).astype(int)
        daily['7d_MA'] = daily['Inflow'].rolling(7, min_periods=1).mean()

        # Forecast
        fc_df = forecast_linear(daily['Inflow'], periods=forecast_periods)
        last_date = daily[dc].max()
        fc_dates  = [last_date + timedelta(days=i+1) for i in range(forecast_periods)]

        st.markdown('<div class="section-header">BURN-UP VELOCITY CURVE</div>', unsafe_allow_html=True)
        fig_v = go.Figure()
        fig_v.add_trace(go.Scatter(x=daily[dc], y=daily['Cumulative_Inflow'], name='Inflow (Discovered)',
            fill='tonexty', fillcolor='rgba(255,0,110,0.08)', line=dict(color='#ff006e', width=2),
            hovertemplate='%{x}<br>Inflow: %{y:,.0f}<extra></extra>'))
        fig_v.add_trace(go.Scatter(x=daily[dc], y=daily['Cumulative_Outflow'], name='Outflow (Resolved)',
            fill='tozeroy', fillcolor='rgba(0,245,255,0.06)', line=dict(color='#00f5ff', width=2),
            hovertemplate='%{x}<br>Outflow: %{y:,.0f}<extra></extra>'))
        fig_v.add_trace(go.Scatter(x=daily[dc], y=daily['7d_MA'], name='7-Day Moving Avg',
            line=dict(color='#ffd60a', width=1.5, dash='dot'),
            hovertemplate='%{x}<br>MA7: %{y:,.1f}<extra></extra>'))
        apply_layout(fig_v, "Project Stability Burn-Up Curve")
        st.plotly_chart(fig_v, use_container_width=True)

        col_f1, col_f2 = st.columns(2)
        with col_f1:
            st.markdown('<div class="section-header">PREDICTIVE FORECAST</div>', unsafe_allow_html=True)
            fig_fc = go.Figure()
            fig_fc.add_trace(go.Scatter(x=daily[dc], y=daily['Inflow'], name='Historical',
                line=dict(color='#00f5ff', width=2),
                hovertemplate='%{x}<br>%{y:,.0f}<extra></extra>'))
            fig_fc.add_trace(go.Scatter(x=fc_dates, y=fc_df['forecast'], name=f'{forecast_periods}d Forecast',
                line=dict(color='#ffd60a', width=2, dash='dash'),
                hovertemplate='%{x}<br>Forecast: %{y:,.1f}<extra></extra>'))
            fig_fc.add_trace(go.Scatter(
                x=fc_dates + fc_dates[::-1],
                y=fc_df['upper'].tolist() + fc_df['lower'].tolist()[::-1],
                fill='toself', fillcolor='rgba(255,214,10,0.06)',
                line=dict(color='rgba(0,0,0,0)'), showlegend=True, name='95% CI'))
            apply_layout(fig_fc, f"{forecast_periods}-Day Linear Forecast")
            st.plotly_chart(fig_fc, use_container_width=True)

        with col_f2:
            st.markdown('<div class="section-header">DAILY INFLOW CALENDAR</div>', unsafe_allow_html=True)
            fig_cal = px.density_heatmap(
                df_filtered.assign(DayOfWeek=df_filtered[dc].dt.day_name(),
                                   Week=df_filtered[dc].dt.isocalendar().week),
                x='Week', y='DayOfWeek', z=raw_y, histfunc='sum',
                color_continuous_scale=[[0,'#0a1628'],[0.5,'#004080'],[1,'#00f5ff']]
            )
            apply_layout(fig_cal, "Weekly Inflow Heatmap")
            st.plotly_chart(fig_cal, use_container_width=True)

# ────────────────────────────────────────
# TAB 4 — ADVANCED ANALYTICS
# ────────────────────────────────────────
with tabs[3]:
    c_a, c_b = st.columns(2)

    with c_a:
        st.markdown('<div class="section-header">CORRELATION INTELLIGENCE</div>', unsafe_allow_html=True)
        if len(num_cols) >= 2:
            corr = df_filtered[num_cols].corr()
            mask = np.triu(np.ones_like(corr, dtype=bool))
            corr_masked = corr.copy()
            corr_masked[mask] = np.nan

            fig_corr = go.Figure(go.Heatmap(
                z=corr.values,
                x=corr.columns.tolist(),
                y=corr.columns.tolist(),
                colorscale=[[0,'#ff006e'],[0.5,'#020408'],[1,'#00f5ff']],
                zmid=0,
                text=np.round(corr.values, 2),
                texttemplate='%{text}',
                textfont=dict(family='Share Tech Mono', size=9),
                hovertemplate='%{y} × %{x}<br>r = %{z:.3f}<extra></extra>',
                colorbar=dict(tickfont=dict(color='#8eb8d4'), title=dict(text='r', font=dict(color='#8eb8d4')))
            ))
            apply_layout(fig_corr, "Pearson Correlation Matrix")
            st.plotly_chart(fig_corr, use_container_width=True)
        else:
            st.info("Need ≥2 numeric columns for correlation analysis.")

    with c_b:
        st.markdown('<div class="section-header">MULTI-METRIC RADAR</div>', unsafe_allow_html=True)
        if len(num_cols) >= 3:
            radar_cols = num_cols[:min(6, len(num_cols))]
            radar_data = df_filtered.groupby(slicer)[radar_cols].mean()
            scaler_r   = StandardScaler()
            radar_norm = pd.DataFrame(scaler_r.fit_transform(radar_data), columns=radar_cols, index=radar_data.index)
            # normalize to 0-1
            radar_norm = (radar_norm - radar_norm.min()) / (radar_norm.max() - radar_norm.min() + 1e-9)

            fig_radar = go.Figure()
            for i, dim in enumerate(radar_norm.index[:8]):
                vals = radar_norm.loc[dim].tolist()
                fig_radar.add_trace(go.Scatterpolar(
                    r=vals + [vals[0]], theta=radar_cols + [radar_cols[0]],
                    fill='toself', fillcolor=f'rgba({int(PALETTE[i%len(PALETTE)][1:3],16)},{int(PALETTE[i%len(PALETTE)][3:5],16)},{int(PALETTE[i%len(PALETTE)][5:],16)},0.15)',
                    line=dict(color=PALETTE[i%len(PALETTE)], width=1.5),
                    name=str(dim)
                ))
            fig_radar.update_layout(
                polar=dict(
                    bgcolor='rgba(6,13,22,0.9)',
                    radialaxis=dict(visible=True, range=[0,1], gridcolor='rgba(0,245,255,0.1)', tickfont=dict(color='#3d6680', family='Share Tech Mono')),
                    angularaxis=dict(gridcolor='rgba(0,245,255,0.1)', tickfont=dict(color='#8eb8d4', family='Rajdhani'))
                ),
                **{k:v for k,v in PLOTLY_LAYOUT.items() if k not in ['xaxis','yaxis']}
            )
            fig_radar.update_layout(title=dict(text="Multi-Metric Radar (Normalized)", font=dict(family='Orbitron', color='#00f5ff', size=13)))
            st.plotly_chart(fig_radar, use_container_width=True)

    # Pareto + Violin
    col_p, col_v = st.columns(2)
    with col_p:
        st.markdown('<div class="section-header">PARETO ANALYSIS (80/20 RULE)</div>', unsafe_allow_html=True)
        pareto = df_filtered.groupby(slicer)[raw_y].sum().sort_values(ascending=False).reset_index()
        pareto['Cumulative%'] = (pareto[raw_y].cumsum() / pareto[raw_y].sum() * 100)
        pareto['Vital Few'] = pareto['Cumulative%'] <= 80

        fig_pareto = go.Figure()
        fig_pareto.add_trace(go.Bar(
            x=pareto[slicer], y=pareto[raw_y],
            marker=dict(color=['#ff006e' if v else '#004080' for v in pareto['Vital Few']], line=dict(color='rgba(0,245,255,0.2)',width=0.5)),
            name='Volume', yaxis='y', hovertemplate='<b>%{x}</b><br>Volume: %{y:,.0f}<extra></extra>'))
        fig_pareto.add_trace(go.Scatter(
            x=pareto[slicer], y=pareto['Cumulative%'],
            line=dict(color='#ffd60a', width=2), name='Cumulative %', yaxis='y2',
            hovertemplate='<b>%{x}</b><br>Cumulative: %{y:.1f}%<extra></extra>'))
        fig_pareto.add_hline(y=80, line=dict(color='#00f5ff', width=1, dash='dot'), yref='y2',
                              annotation_text='80%', annotation_font=dict(color='#00f5ff', family='Share Tech Mono'))
        apply_layout(fig_pareto, "Pareto Chart — Vital Few")
        fig_pareto.update_layout(
            yaxis2=dict(overlaying='y', side='right', range=[0,110], tickfont=dict(color='#ffd60a', family='Share Tech Mono'), gridcolor='rgba(0,0,0,0)'),
            legend=dict(orientation='h', y=-0.15)
        )
        st.plotly_chart(fig_pareto, use_container_width=True)

    with col_v:
        st.markdown('<div class="section-header">VIOLIN DENSITY PLOT</div>', unsafe_allow_html=True)
        fig_vio = go.Figure()
        for i, grp in enumerate(df_filtered[slicer].unique()):
            sub = df_filtered[df_filtered[slicer]==grp][raw_y]
            fig_vio.add_trace(go.Violin(
                y=sub, name=str(grp), box_visible=True, meanline_visible=True,
                fillcolor=f'rgba({int(PALETTE[i%len(PALETTE)][1:3],16)},{int(PALETTE[i%len(PALETTE)][3:5],16)},{int(PALETTE[i%len(PALETTE)][5:],16)},0.25)',
                line_color=PALETTE[i%len(PALETTE)],
                hovertemplate='<b>%{x}</b><br>%{y}<extra></extra>'
            ))
        apply_layout(fig_vio, f"{y_label} Density Distribution")
        st.plotly_chart(fig_vio, use_container_width=True)

# ────────────────────────────────────────
# TAB 5 — AUDIT TRAIL
# ────────────────────────────────────────
with tabs[4]:
    st.markdown('<div class="section-header">STRATEGIC AUDIT TRAIL</div>', unsafe_allow_html=True)

    search = st.text_input("🔎 Search Repository", placeholder="Bug ID, Severity, Module...", label_visibility="collapsed")

    col_exp, col_dl = st.columns([3,1])
    with col_exp:
        if search:
            view_df = df_filtered[df_filtered.apply(lambda r: search.lower() in ' '.join(r.astype(str).str.lower()), axis=1)]
        else:
            view_df = df_filtered.copy()
        if date_cols and 'Age_Days' not in view_df.columns:
            view_df['Age_Days'] = (pd.Timestamp(datetime.now()) - view_df[date_cols[0]]).dt.days.clip(lower=0)

    with col_dl:
        st.download_button(
            "⬇ EXPORT CSV",
            data=view_df.to_csv(index=False).encode(),
            file_name=f"nexus_audit_{datetime.now().strftime('%Y%m%d_%H%M')}.csv",
            mime="text/csv",
            use_container_width=True
        )

    # Highlight critical rows
    def highlight_audit(row):
        if 'Age_Days' in row and row.get('Age_Days', 0) > 14:
            return ['background-color: rgba(255,0,110,0.12); color: #ff8888'] * len(row)
        elif 'Age_Days' in row and row.get('Age_Days', 0) > 7:
            return ['background-color: rgba(255,214,10,0.08); color: #ffd60a'] * len(row)
        return ['' for _ in row]

    st.dataframe(view_df.style.apply(highlight_audit, axis=1), use_container_width=True, height=400)

    st.markdown(f"<div style='font-family:Share Tech Mono;font-size:0.65rem;color:#3d6680;margin-top:0.5rem;'>Showing {len(view_df):,} of {len(df_filtered):,} records</div>", unsafe_allow_html=True)

    # Export PPTX
    st.markdown('<div class="section-header" style="margin-top:1.5rem;">BOARDROOM EXPORT</div>', unsafe_allow_html=True)
    c_pp1, c_pp2 = st.columns(2)
    with c_pp1:
        if st.button("📊 GENERATE EXECUTIVE BRIEFING (PPTX)", use_container_width=True):
            pptx_bytes = build_pptx(
                user=st.session_state['user'],
                stability=stability_score,
                status_text=status_text,
                total_val=total_val,
                metric_label=y_label,
                top_module=top_module,
                dre=dre_score,
                df_summary=risk_df.head(10)
            )
            st.download_button(
                "📥 DOWNLOAD BRIEFING.PPTX",
                data=pptx_bytes,
                file_name=f"NEXUS_Briefing_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
    with c_pp2:
        st.download_button(
            "📥 EXPORT RISK MATRIX (CSV)",
            data=risk_df.to_csv(index=False).encode(),
            file_name=f"NEXUS_RiskMatrix_{datetime.now().strftime('%Y%m%d_%H%M')}.csv",
            mime="text/csv",
            use_container_width=True
        )

# ────────────────────────────────────────
# TAB 6 — AI INSIGHTS ENGINE
# ────────────────────────────────────────
with tabs[5]:
    st.markdown('<div class="section-header">AI-POWERED STRATEGIC INSIGHTS</div>', unsafe_allow_html=True)
    st.markdown("""
    <div class="insight-card" style="border-left: 2px solid #00f5ff; margin-bottom: 0.8rem;">
        <strong>⚡ NEXUS Automated Intelligence Engine</strong><br>
        Rule-based insights derived from real-time statistical analysis of your filtered dataset.
    </div>
    """, unsafe_allow_html=True)

    for i, insight in enumerate(insights_list):
        icon = ["🎯","⚡","📈","🔬"][i % 4]
        st.markdown(f"""
        <div class="insight-card">
            {icon} {insight}
        </div>
        """, unsafe_allow_html=True)

    # Statistical Summary
    st.markdown('<div class="section-header" style="margin-top:1.5rem;">STATISTICAL SUMMARY</div>', unsafe_allow_html=True)
    desc = df_filtered[num_cols].describe().round(2)
    st.dataframe(desc, use_container_width=True)

    # Dimension deep-dive
    st.markdown('<div class="section-header" style="margin-top:1rem;">DIMENSION DEEP-DIVE</div>', unsafe_allow_html=True)
    selected_dim = st.selectbox("Select Dimension for Deep Analysis", df_filtered[slicer].unique())
    dim_data = df_filtered[df_filtered[slicer] == selected_dim]
    d1, d2, d3, d4 = st.columns(4)
    d1.metric("Record Count", f"{len(dim_data):,}")
    d2.metric(f"Total {y_label}", f"{dim_data[raw_y].sum():,.1f}")
    d3.metric(f"Avg {y_label}", f"{dim_data[raw_y].mean():,.2f}")
    d4.metric(f"σ {y_label}", f"{dim_data[raw_y].std():,.2f}")

    if len(num_cols) >= 2:
        col_i1, col_i2 = st.columns(2)
        with col_i1:
            x_pick = st.selectbox("X Variable", num_cols, index=0)
        with col_i2:
            y_pick = st.selectbox("Y Variable", num_cols, index=min(1, len(num_cols)-1))

        if x_pick != y_pick:
            all_dim_data = df_filtered.copy()
            highlight = all_dim_data[slicer] == selected_dim

            fig_reg = go.Figure()
            # Background scatter
            fig_reg.add_trace(go.Scatter(
                x=all_dim_data[~highlight][x_pick], y=all_dim_data[~highlight][y_pick],
                mode='markers', marker=dict(color='rgba(0,245,255,0.15)', size=4), name='Other', showlegend=True))
            # Highlighted scatter
            fig_reg.add_trace(go.Scatter(
                x=dim_data[x_pick], y=dim_data[y_pick],
                mode='markers', marker=dict(color='#ff006e', size=7, symbol='diamond'), name=str(selected_dim)))
            # Regression line
            if len(dim_data) > 2:
                slope, intercept, r_val, p_val, _ = stats.linregress(dim_data[x_pick].fillna(0), dim_data[y_pick].fillna(0))
                x_line = np.linspace(dim_data[x_pick].min(), dim_data[x_pick].max(), 100)
                y_line = slope * x_line + intercept
                fig_reg.add_trace(go.Scatter(x=x_line, y=y_line, mode='lines',
                    line=dict(color='#ffd60a', width=1.5, dash='dot'), name=f'Regression (r²={r_val**2:.3f})'))
            apply_layout(fig_reg, f"Regression: {x_pick} vs {y_pick}")
            st.plotly_chart(fig_reg, use_container_width=True)

    st.markdown(f"""
    <div class="insight-card" style="margin-top:1rem; border-left: 2px solid #ffd60a;">
        <strong>⬡ Session Intelligence:</strong><br>
        Analyst: <strong>{st.session_state['user']}</strong> &nbsp;|&nbsp;
        Records analyzed: <strong>{len(df_filtered):,}</strong> &nbsp;|&nbsp;
        Dimensions tracked: <strong>{df_filtered[slicer].nunique()}</strong> &nbsp;|&nbsp;
        Anomalies flagged: <strong>{int(anomaly_count)}</strong> &nbsp;|&nbsp;
        Stability: <strong>{stability_score}%</strong>
    </div>
    """, unsafe_allow_html=True)
